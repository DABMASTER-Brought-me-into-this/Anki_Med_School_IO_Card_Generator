# logic.py
# Imports
import os
import cv2
import numpy as np
from rapidocr_onnxruntime import RapidOCR
from groq import Groq
import json
import base64
import time
import re
import glob
from PIL import Image
from itertools import combinations
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
from wand.image import Image as WandImage
from wand.color import Color
import csv


# Functions
def io_generate(input_filename, output_filename, results, temp_dir):
    pathway = os.path.join(temp_dir, f"image{input_filename}.png")

    # Allow CV2 to load it
    image = cv2.imread(pathway)

    count = output_filename
    for result in results:
        # Cords
        cords = result[0]

        # Reshaping Cords to Numpy Array
        all_text_points = np.array(cords, dtype=np.int32)
        all_text_points = all_text_points.reshape((-1, 1, 2))

        cv2.fillPoly(image, [all_text_points], color=(255, 0, 0))
        cv2.putText(
            image,  # Image
            str(count),  # Label
            (int(cords[0][0]), int(cords[0][1]) + 20),  # Cords
            cv2.FONT_HERSHEY_SIMPLEX,  # Font
            1.0,  # Font Size
            (255, 255, 255),  # Font Color
            2  # Font Thickness
        )

        cv2.imwrite(os.path.join(temp_dir, f'covered_image{count}.png'), image)
        print("Text covered successfully.")


    original_image = cv2.imread(os.path.join(temp_dir, f'covered_image{count}.png'))
    for result in results:
        # Make a copy of Image
        current_card = original_image.copy()

        # Cords & Text
        cords = result[0]
        text = result[1]

        # Writing the covered text to a text file
        with open(os.path.join(temp_dir, f"answer{count}.txt"), 'a') as file:
            file.write(f"{text}")

        # Reshaping Cords to Numpy Array
        all_text_points = np.array(cords, dtype=np.int32)
        all_text_points = all_text_points.reshape((-1, 1, 2))

        cv2.fillPoly(current_card, [all_text_points], color=(0, 255, 0))
        cv2.imwrite(os.path.join(temp_dir, f'covered_image{count}.png'), current_card)
        print("Text covered successfully.")

        count += 1

    return count


def standard_generate(input_filename, output_filename, slide_text, temp_dir):
    try:
        # Renaming Image to Final Image
        os.rename(os.path.join(temp_dir, f"image{input_filename}.png"),
                  os.path.join(temp_dir, f"covered_image{output_filename}.png"))

        # Creating Answer File
        with open(os.path.join(temp_dir, f"answer{output_filename}.txt"), 'a') as file:
            file.write(f"{slide_text}")

    except Exception as e:
        print(e)

    return output_filename + 1


def groq_preprocessing_tool(similar_image_group, slide_deck_name, groq_api):
    # Connecting To Groq
    client = Groq(
        api_key=groq_api
    )

    def default_keep_all():
        return {os.path.basename(f): "keep" for f in similar_image_group}

    content_payload = [
        {
            "type": "text",
            "text": f"""
            You are a Medical Content Curator.
            Context: The user is converting the lecture "{slide_deck_name}" into flashcards.

            TASK: The user has identified the following {len(similar_image_group)} images as DUPLICATES or variants of the same slide.
            You must decide which image is the 'Best Version' to keep.

            CRITERIA FOR 'keep':
            1. HIGHEST RESOLUTION: Prefer sharp, clear text.
            2. DIAGNOSTIC VALUE: Prefer images with CLEAR labels/arrows pointing to anatomy.
            3. CLEANLINESS: Avoid heavy watermarks if a cleaner version exists.

            DECISION LOGIC:
            - If Image A is Clean and Image B has Labels: Mark B 'keep', Mark A 'delete'.
            - If Image A and Image B are identical content but A is larger/sharper: Mark A 'keep', Mark B 'delete'.
            - If you cannot decide, mark both as 'keep'.

            Return ONLY a valid JSON object mapping filenames to 'keep' or 'delete'.
            Example format:
            {{
                "covered_image1.png": "keep",
                "covered_image2.png": "delete"
            }}
            """
        }
    ]

    valid_filenames = []
    if similar_image_group:
        for file_path in similar_image_group:
            try:
                with open(file_path, "rb") as image_file:
                    encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
                    filename = os.path.basename(file_path)
                    valid_filenames.append(filename)

                    content_payload.append({
                        "type": "text",
                        "text": f"Image Filename: {filename}"
                    })
                    content_payload.append({
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/png;base64,{encoded_string}",
                        },
                    })
            except FileNotFoundError:
                print(f"Warning: File not found {file_path}")

    try:
        chat_completion = client.chat.completions.create(
            messages=[{"role": "user", "content": content_payload}],
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            temperature=0.0,
            response_format={"type": "json_object"}
        )

        response_content = chat_completion.choices[0].message.content
        start_index = response_content.find('{')
        end_index = response_content.rfind('}')

        if start_index != -1 and end_index != -1:
            clean_json = response_content[start_index: end_index + 1]
            return json.loads(clean_json)
        return default_keep_all()

    except Exception as e:
        print(f"Groq Preprocessing Error: {e}")
        return default_keep_all()


def similarity_images(image1, image2):
    try:
        im1 = cv2.imread(image1, cv2.IMREAD_GRAYSCALE)
        im2 = cv2.imread(image2, cv2.IMREAD_GRAYSCALE)
        im1 = cv2.resize(im1, (64, 64))
        im2 = cv2.resize(im2, (64, 64))
        diff = im1.astype(float) - im2.astype(float)
        squared_diff = diff ** 2
        mse = np.mean(squared_diff)
        return mse < 200
    except FileNotFoundError:
        return False


def pre_processing_prune(slide_deck_name, temp_dir, groq_api):
    search_pattern = os.path.join(temp_dir, "image*.png")
    image_files = glob.glob(search_pattern)
    image_files.sort(key=lambda f: int(re.search(r'\d+', os.path.basename(f)).group()))

    for image in image_files:
        bad_card = False
        try:
            with Image.open(image) as img:
                width, height = img.size
                if width < 150 or height < 150: bad_card = True
                if width / height > 5.0 or width / height < 0.2: bad_card = True
        except Exception as e:
            pass

        try:
            if bad_card: os.remove(image)
        except FileNotFoundError:
            pass

    image_files = glob.glob(search_pattern)
    image_files.sort(key=lambda f: int(re.search(r'\d+', os.path.basename(f)).group()))

    orphans_to_delete = []
    batch_size = 4

    for i in range(0, len(image_files), batch_size):
        batch = image_files[i: i + batch_size]
        combinations_batch = list(combinations(batch, 2))
        similar_images = []
        for combination in combinations_batch:
            if similarity_images(combination[0], combination[1]):
                if combination[0] not in similar_images: similar_images.append(combination[0])
                if combination[1] not in similar_images: similar_images.append(combination[1])

        if len(similar_images) >= 2:
            decisions = groq_preprocessing_tool(similar_images, slide_deck_name, groq_api)
            for key, value in decisions.items():
                if value == "delete": orphans_to_delete.append(key)

    for orphan in orphans_to_delete:
        try:
            os.remove(os.path.join(temp_dir, orphan))
        except Exception as e:
            pass


def post_processing_prune(slide_deck_name, temp_dir, groq_api):
    search_pattern = os.path.join(temp_dir, "answer*.txt")
    answer_files = glob.glob(search_pattern)
    answer_files.sort(key=lambda f: int(re.search(r'\d+', os.path.basename(f)).group()))

    img_pattern = os.path.join(temp_dir, "covered_image*.png")
    covered_image_files = glob.glob(img_pattern)
    covered_image_files.sort(key=lambda f: int(re.search(r'\d+', os.path.basename(f)).group()))

    for file in answer_files:
        bad_card = False
        try:
            with open(file, "r", encoding="utf-8") as txt_file:
                keywords = ["learning objectives", "objectives", "goals", "agenda", "overview", "course outline",
                            "schedule", "announcements", "housekeeping", "welcome to", "module", "lecture",
                            "attendance", "any questions", "questions?", "questions & answers", "q&a", "thank you",
                            "thanks", "contact info", "contact me", "email:", "discussion", "group activity",
                            "quiz time", "pop quiz", "break time", "references", "works cited", "bibliography",
                            "sources", "citations", "image credits", "image source", "figure source", "copyright",
                            "all rights reserved", "conflict of interest", "disclosures", "financial disclosure",
                            "no disclosures", "summary", "conclusion", "recap", "take home points", "key takeaways",
                            "take-away", "in summary", "to summarize", "click to add title", "click to add text",
                            "slide number"]
                content = txt_file.read().lower()
                for keyword in keywords:
                    if keyword in content: bad_card = True

            if bad_card:
                os.remove(file)
                match = re.search(r'answer(\d+)\.txt', os.path.basename(file))
                if match:
                    os.remove(os.path.join(temp_dir, f"covered_image{match.group(1)}.png"))
        except Exception as e:
            pass

    final_image_list = glob.glob(img_pattern)
    pruning_decisions = {}
    batch_size = 5

    for i in range(0, len(final_image_list), batch_size):
        batch = final_image_list[i: i + batch_size]
        print(f" -> Pruning Batch {i // batch_size + 1}: Checking {len(batch)} images...")
        batch_results = groq_pruning_tool(batch, slide_deck_name, groq_api)
        pruning_decisions.update(batch_results)

    for key, value in pruning_decisions.items():
        if value == "Delete":
            try:
                os.remove(os.path.join(temp_dir, key))
                match = re.search(r'\d+', key)
                if match:
                    os.remove(os.path.join(temp_dir, f"answer{match.group()}.txt"))
            except FileNotFoundError:
                pass


def groq_pruning_tool(file_paths, slide_deck_name, groq_api):
    client = Groq(api_key=groq_api)
    content_payload = [
        {
            "type": "text",
            "text": f"""
                You are a rigorous Medical Slide Deck Cleaner. 
                Context: These images are from a lecture deck named "{slide_deck_name}".

                Your Task: Review each image and assign it one of two statuses: 'Delete' or 'Pass'.

                CRITERIA FOR 'Delete' (Be aggressive):
                1. Irrelevant: Comic strips, cartoons, "Questions?" slides, "Thank You" slides, or generic stock photos of people.
                2. Duplicates: If two images IN THIS LIST are effectively the same content, mark the lower-quality one as Delete.
                3. Nonsense: The image is too blurry to read, contains only a slide number/header, or is a corrupted visual artifact.
                4. Formatting Artifacts: Images that are just lines, background textures, or company logos.

                CRITERIA FOR 'Pass':
                1. Medical Content: Contains anatomy, histology, charts, text definitions, or clinical photos.

                Return ONLY a valid JSON object.
                Example format:
                {{
                    "image1.png": "Pass",
                    "image2.png": "Delete",
                    "image3.png": "Delete"
                }}
                """
        }
    ]

    valid_filenames = []
    if file_paths:
        for file_path in file_paths:
            try:
                with open(file_path, "rb") as image_file:
                    encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
                    filename = os.path.basename(file_path)
                    valid_filenames.append(filename)
                    content_payload.append({"type": "text", "text": f"Image Filename: {filename}"})
                    content_payload.append(
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{encoded_string}"}})
            except FileNotFoundError:
                pass

    try:
        chat_completion = client.chat.completions.create(
            messages=[{"role": "user", "content": content_payload}],
            model="meta-llama/llama-4-maverick-17b-128e-instruct",
            temperature=0.0,
            response_format={"type": "json_object"}
        )

        response_content = chat_completion.choices[0].message.content
        start_index = response_content.find('{')
        end_index = response_content.rfind('}')

        if start_index != -1 and end_index != -1:
            clean_json_string = response_content[start_index: end_index + 1]
            try:
                return json.loads(clean_json_string)
            except json.JSONDecodeError:
                return {name: "Pass" for name in valid_filenames}
        else:
            return {name: "Pass" for name in valid_filenames}
    except Exception as e:
        return {name: "Pass" for name in valid_filenames}


# Slide to Text
def run_pipeline(presentation_name, start_counter, temp_dir, groq_api):
    print("Starting Code")
    start_time = time.perf_counter()
    prs = pptx.Presentation(os.path.join(temp_dir, presentation_name))
    slides = prs.slides

    counter = 0
    slide_texts = {}

    print("Reading Slides")
    for slide in slides:
        print("Finding Shapes in new Slide")
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print("Found Image")
                im = shape.image
                im_byte = im.blob
                im_filename = os.path.join(temp_dir, f'image{counter}.png')
                im_stream = io.BytesIO(im_byte)

                print("Storing Image")
                try:
                    if 'wmf' in im.content_type or 'emf' in im.content_type:
                        with WandImage(blob=im_byte) as vector_img:
                            vector_img.background_color = Color('white')
                            vector_img.format = 'png'
                            vector_img.save(filename=im_filename)
                    else:
                        img = Image.open(im_stream)
                        img.save(im_filename, format='PNG')

                    print("Storing Text")
                    slide_text = ""
                    for shape1 in slide.shapes:
                        if shape1.has_text_frame:
                            slide_text += str(shape1.text_frame.text)
                    slide_texts[im_filename] = slide_text

                except Exception as e:
                    print("Error 222: Bad File Type")
                counter += 1

    print("Pre-processing")
    pre_processing_prune(presentation_name, temp_dir, groq_api)

    search_pattern = os.path.join(temp_dir, "image*.png")
    image_list = glob.glob(search_pattern)
    image_list.sort(key=lambda f: int(re.search(r'\d+', os.path.basename(f)).group()))

    print("Loading OCR")
    ocr_engine = RapidOCR()
    counter = int(start_counter)

    print("Running OCR on All Images")
    for image in image_list:
        results, elapse = ocr_engine(image)
        is_valid_io_card = False

        if results:
            for detection in results:
                text = detection[1]
                confidence = detection[2]
                if re.search('[a-zA-Z]', text) and confidence > 0.6:
                    is_valid_io_card = True
                    break

        if is_valid_io_card:
            print("Creating IO Card...")
            num_str = re.search(r'\d+', os.path.basename(image)).group()
            counter = io_generate(int(num_str), counter, results, temp_dir)
        else:
            print("Creating Standard Card...")
            num_str = re.search(r'\d+', os.path.basename(image)).group()
            counter = standard_generate(int(num_str), counter, slide_texts.get(image, ""), temp_dir)

    files = glob.glob(search_pattern)
    print("Pruning Orphan Files")
    for file_path in files:
        if not os.path.basename(file_path).startswith("covered_"):
            try:
                os.remove(file_path)
                print(f"Deleted: {file_path}")
            except Exception as e:
                pass

    print("Post Processing")
    post_processing_prune(presentation_name, temp_dir, groq_api)
    print(f"The slide_to_text took {time.perf_counter() - start_time} seconds.")


# Formatting
def create_csv_file(csv_name, temp_dir):
    print("Formatting")
    search_pattern = os.path.join(temp_dir, "answer*.txt")
    answer_files = glob.glob(search_pattern)
    answer_files.sort(key=lambda f: int(re.search(r'\d+', os.path.basename(f)).group()))

    img_pattern = os.path.join(temp_dir, "covered_image*.png")
    image_files = glob.glob(img_pattern)
    image_files.sort(key=lambda f: int(re.search(r'\d+', os.path.basename(f)).group()))

    answer_ids = set()
    for f in answer_files:
        match = re.search(r"answer(\d+)\.txt", os.path.basename(f))
        if match: answer_ids.add(int(match.group(1)))

    image_ids = set()
    for f in image_files:
        match = re.search(r"covered_image(\d+)\.png", os.path.basename(f))
        if match: image_ids.add(int(match.group(1)))

    valid_ids = sorted(list(answer_ids.intersection(image_ids)))
    csv_rows = []

    for num in valid_ids:
        txt_filename = os.path.join(temp_dir, f"answer{num}.txt")
        img_filename = f"covered_image{num}.png"
        try:
            with open(txt_filename, "r", encoding="utf-8") as file:
                content = file.read().strip().replace("\n", "<br>")
                csv_rows.append([f'<img src="{img_filename}">', content])
        except Exception as e:
            pass

    if not csv_name.endswith(".csv"): csv_name += ".csv"

    if csv_rows:
        with open(os.path.join(temp_dir, csv_name), 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)
            writer.writerows(csv_rows)
        return csv_name
    return None
